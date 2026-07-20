"""AXIOM report artifacts (Phase 7f): themed PPTX board deck + PDF, rendered
from the same board_report() engine payload the on-screen report uses.

AXIOM dark theme — deep pine backgrounds, ivory text, brass accents — mirroring
the brochure/PDF language. Charts are matplotlib PNGs, theme-matched.
"""
import io
from datetime import datetime

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Polygon
import numpy as np

# ---- palette ---------------------------------------------------------------
PINE = "#0E1F17"     # deep pine background
PANEL = "#15291F"    # raised panel
IVORY = "#F4F1E8"    # primary text
BRASS = "#C9A24B"    # accent
GREEN = "#4ADE80"    # positive
RED = "#F0776B"      # negative
AMBER = "#FBBF24"
MUTED = "#9FB3A8"    # secondary text
GRID = "#26402F"

_RGB = lambda h: tuple(int(h[i:i + 2], 16) for i in (1, 3, 5))


def issued_line(issued_at: datetime, dataset_version) -> str:
    dv = f" · Data version {dataset_version}" if dataset_version is not None else ""
    return f"Issued {issued_at.strftime('%d %b %Y, %H:%M')}{dv}"


def report_filename(company: str, report_type: str, fmt: str, issued_at: datetime) -> str:
    safe = "".join(ch if (ch.isalnum() or ch in " -") else "" for ch in (company or "Company")).strip()
    safe = "_".join(safe.split()) or "Company"
    rt = "_".join((report_type or "Report").split())
    return f"{safe}_{rt}_{issued_at.strftime('%Y-%m-%d_%H%M')}.{fmt}"


# ---- matplotlib helpers ----------------------------------------------------
def _fig(w=8.0, h=4.2):
    fig, ax = plt.subplots(figsize=(w, h), dpi=150)
    fig.patch.set_facecolor(PINE)
    ax.set_facecolor(PINE)
    for s in ax.spines.values():
        s.set_color(GRID)
    ax.tick_params(colors=MUTED, labelsize=8)
    ax.title.set_color(IVORY)
    ax.xaxis.label.set_color(MUTED); ax.yaxis.label.set_color(MUTED)
    return fig, ax


def _png(fig) -> bytes:
    buf = io.BytesIO()
    fig.tight_layout(pad=1.1)
    fig.savefig(buf, format="png", facecolor=PINE, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def chart_fan(fan, title="Revenue — projection with P05–P95 band"):
    if not fan:
        return None
    yrs = [p["year"] for p in fan]
    p05 = [p.get("p05") for p in fan]; p25 = [p.get("p25") for p in fan]
    p50 = [p.get("p50") for p in fan]; p75 = [p.get("p75") for p in fan]
    p95 = [p.get("p95") for p in fan]
    fig, ax = _fig()
    ax.fill_between(yrs, p05, p95, color=BRASS, alpha=0.16, label="P05–P95")
    ax.fill_between(yrs, p25, p75, color=BRASS, alpha=0.34, label="P25–P75")
    ax.plot(yrs, p50, color=BRASS, lw=2.4, marker="o", ms=4, label="Median")
    ax.set_title(title, loc="left", fontsize=11, fontweight="bold")
    ax.set_xticks(yrs)
    leg = ax.legend(loc="upper left", fontsize=8, frameon=False)
    for t in leg.get_texts(): t.set_color(MUTED)
    ax.grid(True, color=GRID, lw=0.5, alpha=0.6)
    return _png(fig)


def chart_valuation_lenses(dcf, real_options):
    labels, vals, cols = [], [], []
    if dcf:
        labels.append("DCF (EV)"); vals.append(dcf.get("enterprise_value")); cols.append(BRASS)
        mc = dcf.get("monte_carlo") or {}
        if mc.get("mean") is not None:
            labels.append("Monte Carlo\nmean"); vals.append(mc["mean"]); cols.append(GREEN)
        if mc.get("cvar95") is not None:
            labels.append("MC 95%\ntail (CVaR)"); vals.append(mc["cvar95"]); cols.append(AMBER)
    try:
        rov = real_options["options"]["expand"]["flexibility_value"]
        labels.append("Real-option\nexpand"); vals.append(rov); cols.append(MUTED)
    except (KeyError, TypeError):
        pass
    labels = [l for l, v in zip(labels, vals) if v is not None]
    cols = [c for c, v in zip(cols, vals) if v is not None]
    vals = [v for v in vals if v is not None]
    if not vals:
        return None
    fig, ax = _fig(7.6, 4.0)
    bars = ax.bar(range(len(vals)), vals, color=cols, width=0.6)
    ax.set_xticks(range(len(vals))); ax.set_xticklabels(labels, fontsize=8, color=MUTED)
    ax.set_title("Valuation — three independent lenses", loc="left", fontsize=11, fontweight="bold")
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v, f"{v:,.0f}", ha="center", va="bottom",
                color=IVORY, fontsize=8)
    ax.grid(True, axis="y", color=GRID, lw=0.5, alpha=0.6)
    return _png(fig)


def chart_tornado(recs):
    items = [(r.get("title", r.get("move", "")), r.get("expected_ev_impact"))
             for r in (recs or []) if r.get("expected_ev_impact") is not None]
    if not items:
        return None
    items.sort(key=lambda t: abs(t[1]))
    labels = [t[0][:38] for t in items]; vals = [t[1] for t in items]
    cols = [GREEN if v >= 0 else RED for v in vals]
    fig, ax = _fig(8.0, 3.2 + 0.35 * len(items))
    ax.barh(range(len(vals)), vals, color=cols, height=0.6)
    ax.set_yticks(range(len(vals))); ax.set_yticklabels(labels, fontsize=8, color=IVORY)
    ax.axvline(0, color=MUTED, lw=0.8)
    ax.set_title("Value drivers — EV impact of each lever ($M)", loc="left", fontsize=11, fontweight="bold")
    for i, v in enumerate(vals):
        ax.text(v, i, f" {v:+,.1f}", va="center", ha="left" if v >= 0 else "right",
                color=IVORY, fontsize=8)
    ax.grid(True, axis="x", color=GRID, lw=0.5, alpha=0.6)
    return _png(fig)


def chart_frontier(points, recommended):
    if not points:
        return None
    x = [p["safety_tail_margin"] for p in points]; y = [p["value_mean_ev"] for p in points]
    fig, ax = _fig(7.6, 4.0)
    ax.plot(x, y, color=BRASS, lw=1.6, marker="o", ms=3, alpha=0.8)
    if recommended:
        ax.scatter([recommended["safety_tail_margin"]], [recommended["value_mean_ev"]],
                   s=160, facecolor="none", edgecolor=GREEN, lw=2.2, zorder=5, label="recommended")
        leg = ax.legend(loc="lower right", fontsize=8, frameon=False)
        for t in leg.get_texts(): t.set_color(MUTED)
    ax.set_title("Value–risk frontier over capital structure", loc="left", fontsize=11, fontweight="bold")
    ax.set_xlabel("safety (tail margin)"); ax.set_ylabel("mean EV")
    ax.grid(True, color=GRID, lw=0.5, alpha=0.6)
    return _png(fig)


def chart_radar(l1_subscores):
    pts = [(o.get("title", o.get("code", "")), o.get("score")) for o in (l1_subscores or [])
           if o.get("score") is not None]
    if len(pts) < 3:
        return None
    labels = [p[0][:16] for p in pts]; vals = [p[1] for p in pts]
    n = len(vals); ang = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
    vals2 = vals + vals[:1]; ang2 = ang + ang[:1]
    fig = plt.figure(figsize=(5.4, 5.4), dpi=150); fig.patch.set_facecolor(PINE)
    ax = fig.add_subplot(111, polar=True); ax.set_facecolor(PINE)
    ax.plot(ang2, vals2, color=BRASS, lw=2)
    ax.fill(ang2, vals2, color=BRASS, alpha=0.25)
    ax.set_xticks(ang); ax.set_xticklabels(labels, fontsize=7, color=MUTED)
    ax.set_ylim(0, 10); ax.set_yticks([2, 4, 6, 8, 10])
    ax.tick_params(colors=MUTED); ax.grid(color=GRID)
    ax.spines["polar"].set_color(GRID)
    ax.set_title("Composite Excellence Index — L1 subscores", color=IVORY, fontsize=11,
                 fontweight="bold", pad=18)
    return _png(fig)


def chart_trend(points, title, ylabel=""):
    """points: [{year, value, kind:'hist'|'fcst'}] — solid history, dashed forecast."""
    pts = [p for p in (points or []) if p.get("value") is not None]
    if len(pts) < 2:
        return None
    fig, ax = _fig(7.8, 3.6)
    hx = [p["year"] for p in pts if p.get("kind") != "fcst"]
    hy = [p["value"] for p in pts if p.get("kind") != "fcst"]
    fx = [p["year"] for p in pts if p.get("kind") == "fcst"]
    fy = [p["value"] for p in pts if p.get("kind") == "fcst"]
    if hx:
        ax.plot(hx, hy, color=BRASS, lw=2.4, marker="o", ms=4, label="actual")
    if fx:
        bx = ([hx[-1]] + fx) if hx else fx
        by = ([hy[-1]] + fy) if hx else fy
        ax.plot(bx, by, color=BRASS, lw=2.0, ls="--", marker="o", ms=4, alpha=0.85, label="forecast")
    ax.set_title(title, loc="left", fontsize=11, fontweight="bold")
    ax.set_ylabel(ylabel)
    if hx and fx:
        leg = ax.legend(loc="upper left", fontsize=8, frameon=False)
        for t in leg.get_texts(): t.set_color(MUTED)
    ax.grid(True, color=GRID, lw=0.5, alpha=0.6)
    return _png(fig)


def chart_mc_hist(hist, percentiles=None):
    if not hist or not hist.get("counts"):
        return None
    counts = hist["counts"]; start = hist.get("bin_start", 0); width = hist.get("bin_width", 1)
    xs = [start + i * width for i in range(len(counts))]
    fig, ax = _fig(7.8, 3.8)
    ax.bar(xs, counts, width=width * 0.92, color=BRASS, alpha=0.75)
    for key, col in (("p05", RED), ("p50", IVORY), ("p95", GREEN)):
        v = (percentiles or {}).get(key)
        if v is not None:
            ax.axvline(v, color=col, lw=1.4, ls="--")
            ax.text(v, max(counts) * 0.96, key, color=col, fontsize=7, ha="center")
    ax.set_title("Enterprise value — Monte Carlo distribution", loc="left", fontsize=11, fontweight="bold")
    ax.grid(True, axis="y", color=GRID, lw=0.5, alpha=0.5)
    return _png(fig)


def chart_bars(labels, values, title, colors=None, horizontal=False, fmt="{:,.1f}"):
    vals = [(l, v) for l, v in zip(labels or [], values or []) if v is not None]
    if not vals:
        return None
    labels = [v[0] for v in vals]; values = [v[1] for v in vals]
    colors = colors or [BRASS] * len(values)
    fig, ax = _fig(7.8, max(3.0, 0.5 * len(values) + 1.6) if horizontal else 3.8)
    if horizontal:
        ax.barh(range(len(values)), values, color=colors, height=0.62)
        ax.set_yticks(range(len(values))); ax.set_yticklabels(labels, fontsize=8, color=IVORY)
        for i, v in enumerate(values):
            ax.text(v, i, " " + fmt.format(v), va="center", color=IVORY, fontsize=8)
    else:
        ax.bar(range(len(values)), values, color=colors, width=0.62)
        ax.set_xticks(range(len(values))); ax.set_xticklabels(labels, fontsize=8, color=MUTED, rotation=0)
        for i, v in enumerate(values):
            ax.text(i, v, fmt.format(v), ha="center", va="bottom", color=IVORY, fontsize=8)
    ax.set_title(title, loc="left", fontsize=11, fontweight="bold")
    ax.grid(True, axis=("x" if horizontal else "y"), color=GRID, lw=0.5, alpha=0.55)
    return _png(fig)


def chart_grouped_bars(labels, series_a, series_b, name_a, name_b, title):
    """Two-series grouped bars (e.g. company vs sector)."""
    idx = [i for i, (a, b) in enumerate(zip(series_a, series_b)) if a is not None and b is not None]
    if not idx:
        return None
    labels = [labels[i] for i in idx]; a = [series_a[i] for i in idx]; b = [series_b[i] for i in idx]
    x = np.arange(len(labels)); w = 0.38
    fig, ax = _fig(8.0, 3.8)
    ax.bar(x - w / 2, a, w, color=BRASS, label=name_a)
    ax.bar(x + w / 2, b, w, color=MUTED, label=name_b)
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=8, color=MUTED)
    ax.set_title(title, loc="left", fontsize=11, fontweight="bold")
    leg = ax.legend(loc="upper right", fontsize=8, frameon=False)
    for t in leg.get_texts(): t.set_color(MUTED)
    ax.grid(True, axis="y", color=GRID, lw=0.5, alpha=0.55)
    return _png(fig)


def chart_paths(sample_paths, years, title, color=BRASS, cap=40):
    if not sample_paths or not years:
        return None
    fig, ax = _fig(7.8, 3.8)
    for path in sample_paths[:cap]:
        if len(path) == len(years):
            ax.plot(years, path, color=color, lw=0.6, alpha=0.22)
    ax.set_title(title, loc="left", fontsize=11, fontweight="bold")
    ax.set_xticks(years)
    ax.grid(True, color=GRID, lw=0.5, alpha=0.5)
    return _png(fig)


from PIL import Image, ImageDraw, ImageFont
import matplotlib.font_manager as _fm


def _px_size(png: bytes):
    try:
        return Image.open(io.BytesIO(png)).size      # (w, h)
    except Exception:
        return None


def _fit_box(iw, ih, max_w, max_h):
    """Fit (iw,ih) inside (max_w,max_h) preserving aspect — never stretch."""
    if not iw or not ih:
        return max_w, max_h
    scale = min(max_w / iw, max_h / ih)
    return iw * scale, ih * scale


def wordmark_png(text, bg_hex=PINE, fg_hex=IVORY, accent_hex=BRASS, w=680, h=260):
    """A simple, tasteful text wordmark (for the showcase companies)."""
    img = Image.new("RGB", (w, h), _RGB(bg_hex))
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype(_fm.findfont("DejaVu Sans:bold"), 62)
    except Exception:
        font = ImageFont.load_default()
    bb = d.textbbox((0, 0), text, font=font)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    x, y = (w - tw) / 2 - bb[0], (h - th) / 2 - bb[1]
    d.text((x, y), text, fill=_RGB(fg_hex), font=font)
    d.rectangle([(w - tw) / 2, y + th + 18, (w - tw) / 2 + tw, y + th + 26], fill=_RGB(accent_hex))
    buf = io.BytesIO(); img.save(buf, "PNG")
    return buf.getvalue()


def _logo_from_meta(meta):
    lg = (meta or {}).get("logo")
    if not lg:
        return None
    png = lg[0] if isinstance(lg, (tuple, list)) else lg.get("bytes")
    return png or None


# ============================================================================
# PPTX board deck
# ============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

_C = lambda h: RGBColor.from_string(h.lstrip("#"))
EMU_IN = 914400
SW, SH = 13.333, 7.5


def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _C(PINE)
    return s


def _tb(slide, l, t, w, h, text, size=14, color=IVORY, bold=False, align=PP_ALIGN.LEFT,
        italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = _C(color); r.font.name = "Georgia"
    return box


def _rule(slide, l, t, w, color=BRASS, h=0.035):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _C(color); sh.line.fill.background()
    return sh


def _panel(slide, l, t, w, h, color=PANEL):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _C(color); sh.line.color.rgb = _C(GRID); sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _header(slide, kicker, title, logo_png=None):
    _tb(slide, 0.6, 0.42, 10.5, 0.34, kicker.upper(), size=11, color=BRASS, bold=True)
    _tb(slide, 0.6, 0.74, 10.6, 0.9, title, size=26, color=IVORY, bold=True)
    _rule(slide, 0.62, 1.62, 3.2)
    if logo_png:                                     # small client mark, top-right corner
        sz = _px_size(logo_png)
        w, h = _fit_box(sz[0], sz[1], 1.25, 0.55) if sz else (1.1, 0.5)
        try:
            slide.shapes.add_picture(io.BytesIO(logo_png), Inches(13.333 - 0.5 - w),
                                     Inches(0.42), Inches(w), Inches(h))
        except Exception:
            pass


def _pic(slide, png, l, t, w, h):
    if not png:
        return
    slide.shapes.add_picture(io.BytesIO(png), Inches(l), Inches(t), Inches(w), Inches(h))


def _table(slide, headers, rows, l, t, w, col_w=None, fsize=11, row_h=0.34):
    nrows, ncols = len(rows) + 1, len(headers)
    h = row_h * nrows
    gt = slide.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Inches(cw)
    for j, htext in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = _C(BRASS)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = str(htext); r.font.size = Pt(fsize); r.font.bold = True
        r.font.color.rgb = _C(PINE); r.font.name = "Georgia"
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = _C(PANEL if i % 2 else PINE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = "" if val is None else str(val)
            r.font.size = Pt(fsize); r.font.color.rgb = _C(IVORY); r.font.name = "Georgia"
    return gt


def _rag_word(rag):
    return {"green": "Green", "amber": "Amber", "red": "Red"}.get(rag, "—")


def _fmt(v, pct=False):
    if v is None:
        return "—"
    return f"{v*100:.1f}%" if pct else f"{v:,.1f}"


def build_pptx(report: dict, extras: dict, meta: dict) -> bytes:
    extras = extras or {}
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    S = {s["id"]: s for s in report.get("sections", [])}
    company = report.get("company", {})
    cur = company.get("currency", "")
    issued = meta["issued_at"]
    logo_png = _logo_from_meta(meta)

    # 1. Title — the client's logo leads (it's their board deck); AXIOM subordinate
    s = _blank(prs)
    if logo_png:
        sz = _px_size(logo_png)
        w, h = _fit_box(sz[0], sz[1], 5.2, 1.9) if sz else (4.4, 1.6)
        try:
            s.shapes.add_picture(io.BytesIO(logo_png), Inches(0.9), Inches(1.35),
                                 Inches(w), Inches(h))
        except Exception:
            logo_png = None
    if not logo_png:
        _tb(s, 0.9, 2.15, 11.5, 0.5, "AXIOM", size=18, color=BRASS, bold=True)
    top = 3.55 if _logo_from_meta(meta) else 2.75
    _tb(s, 0.9, top, 11.5, 1.4, meta.get("company_name", company.get("name", "Company")),
        size=40, color=IVORY, bold=True)
    _tb(s, 0.9, top + 1.35, 11.5, 0.6, meta.get("report_type", "Board Report"), size=22, color=MUTED)
    _rule(s, 0.94, top + 2.15, 4.0)
    _tb(s, 0.9, top + 2.35, 11.5, 0.5, issued_line(issued, meta.get("dataset_version")),
        size=13, color=MUTED)
    _tb(s, 0.9, 6.9, 11.5, 0.4, "Confidential — prepared by AXIOM Dynamics", size=10,
        color=MUTED, italic=True)

    # 2. Executive summary
    s = _blank(prs); _header(s, "Executive summary", "The company at a glance", logo_png)
    sm = S.get("summary", {}); sc = sm.get("scorecard", {}); hm = sm.get("headline_metric", {})
    _panel(s, 0.6, 1.95, 5.4, 2.0)
    _tb(s, 0.85, 2.15, 5.0, 0.4, hm.get("label", report.get("headline", {}).get("label", "Enterprise Value")),
        size=13, color=MUTED)
    _tb(s, 0.85, 2.55, 5.0, 0.9, f"{hm.get('value', report.get('headline', {}).get('value', 0)):,.0f} {cur}",
        size=34, color=BRASS, bold=True)
    _tb(s, 0.85, 3.5, 5.0, 0.35, sm.get("takeaway", "")[:120], size=11, color=MUTED, italic=True)
    cards = [("Health index", _fmt(sc.get("health_index"))),
             ("Risk grade", sc.get("risk_grade", "—")),
             ("Optimization", sc.get("optimization_status", "—")),
             ("Optimizer uplift", f"{sc.get('optimization_uplift', 0):,.0f} {cur}")]
    for i, (k, v) in enumerate(cards):
        cx = 6.35 + (i % 2) * 3.35; cy = 1.95 + (i // 2) * 1.05
        _panel(s, cx, cy, 3.15, 0.92)
        _tb(s, cx + 0.2, cy + 0.12, 2.8, 0.3, k.upper(), size=9, color=MUTED, bold=True)
        _tb(s, cx + 0.2, cy + 0.4, 2.8, 0.45, str(v), size=15, color=IVORY, bold=True)
    recs = extras.get("recommendations") or []
    n_pos = sum(1 for r in recs if (r.get("expected_ev_impact") or 0) > 0)
    _tb(s, 0.6, 4.2, 12, 2.6, [
        "The four answers:",
        *[f"•  {w}" for w in sm.get("four_answers", [])[:4]],
        f"Value-creating levers identified: {n_pos}",
    ], size=12, color=IVORY)

    # 3. Statement highlights (KPIs)
    s = _blank(prs); _header(s, "Financial highlights", "Latest actuals vs. prior period", logo_png)
    kpis = S.get("diagnostic", {}).get("kpi_strip") or []
    rows = [[k.get("kpi"), f"{k.get('current'):,.1f}" if k.get("current") is not None else "—",
             f"{k.get('previous'):,.1f}" if k.get("previous") is not None else "—",
             _fmt(k.get("trend"), pct=True)] for k in kpis[:8]]
    if rows:
        _table(s, ["KPI", "Current", "Prior", "Trend"], rows, 0.6, 2.0, 8.5,
               col_w=[3.4, 1.7, 1.7, 1.7])
    _tb(s, 9.4, 2.0, 3.4, 4.6, "Figures in " + (report.get("units_note", "$ millions")),
        size=10, color=MUTED, italic=True)

    # 4. Statement highlights — forecast bands (P05–P95)
    s = _blank(prs); _header(s, "Forecast bands", "Revenue distribution by year (P05–P95)", logo_png)
    fan = (S.get("outlook", {}).get("simulation_baseline", {}) or {}).get("revenue_fan") or []
    rows = [[p.get("year"), f"{p.get('p05'):,.0f}", f"{p.get('p25'):,.0f}",
             f"{p.get('p50'):,.0f}", f"{p.get('p75'):,.0f}", f"{p.get('p95'):,.0f}"] for p in fan]
    if rows:
        _table(s, ["Year", "P05", "P25", "Median", "P75", "P95"], rows, 0.6, 2.1, 11.0)
    else:
        _tb(s, 0.6, 2.3, 11, 0.6, "No probabilistic forecast available for this dataset.",
            size=13, color=MUTED)

    # 5. Forecast fan
    s = _blank(prs); _header(s, "Outlook", "Probabilistic revenue forecast", logo_png)
    _pic(s, chart_fan(fan), 0.7, 1.95, 11.9, 5.0)

    # 6. Valuation lenses
    s = _blank(prs); _header(s, "Valuation", "Three independent lenses", logo_png)
    va = S.get("valuation", {})
    _pic(s, chart_valuation_lenses(va.get("dcf"), va.get("real_options")), 0.9, 1.95, 8.0, 4.7)
    dcf = va.get("dcf", {}); mc = dcf.get("monte_carlo", {})
    _panel(s, 9.2, 2.0, 3.5, 4.5)
    _tb(s, 9.42, 2.2, 3.1, 4.2, [
        "DCF enterprise value",
        f"{dcf.get('enterprise_value', 0):,.0f} {cur}", "",
        "Monte Carlo mean",
        f"{mc.get('mean', 0):,.0f} {cur}", "",
        "95% tail (CVaR)",
        f"{mc.get('cvar95', 0):,.0f} {cur}", "",
        f"WACC {(_fmt(dcf.get('wacc'), pct=True))}",
    ], size=12, color=IVORY)

    # 7. Value drivers / tornado
    s = _blank(prs); _header(s, "Value drivers", "EV impact of each lever", logo_png)
    _pic(s, chart_tornado(S.get("actions", {}).get("recommendations")), 0.7, 1.9, 11.9, 5.1)

    # 8. Risk indicators
    s = _blank(prs); _header(s, "Risk", "Key risk indicators", logo_png)
    rg = S.get("diagnostic", {}).get("risk_grade", {})
    runway = S.get("outlook", {}).get("cash_runway", {})
    cov = S.get("outlook", {}).get("coverage", {})
    cards = [("Risk grade", rg.get("grade", sm.get("scorecard", {}).get("risk_grade", "—"))),
             ("P(cash < 0) ever", _fmt(runway.get("p_cash_below_zero_ever"), pct=True)),
             ("Months to zero", (f"{runway.get('deterministic_months_to_zero'):,.0f}"
                                 if runway.get("deterministic_months_to_zero") is not None else "—")),
             ("Burning cash", "Yes" if runway.get("burning_cash") else "No")]
    for i, (k, v) in enumerate(cards):
        cx = 0.6 + (i % 4) * 3.15
        _panel(s, cx, 2.1, 2.95, 1.4)
        _tb(s, cx + 0.2, 2.32, 2.6, 0.4, k.upper(), size=10, color=MUTED, bold=True)
        _tb(s, cx + 0.2, 2.78, 2.6, 0.55, str(v), size=18, color=IVORY, bold=True)
    _tb(s, 0.6, 3.9, 12, 2.6, [f"•  {w}" for w in (S.get("outlook", {}).get("narrative") or [])[:4]],
        size=12, color=IVORY)

    # 9. CEI + SWOT
    s = _blank(prs); _header(s, "Organizational excellence", "Assessment & SWOT", logo_png)
    cei = extras.get("cei"); swot = extras.get("swot")
    if cei and cei.get("cei") is not None:
        _pic(s, chart_radar(cei.get("l1_subscores")), 0.7, 1.85, 5.2, 5.2)
        _tb(s, 6.2, 2.0, 6.4, 0.6, f"CEI  {cei.get('cei'):.1f} / 10", size=22, color=BRASS, bold=True)
        if swot and swot.get("has_data"):
            cnt = swot.get("counts", {})
            rows = [["Strengths", cnt.get("strengths", 0)], ["Weaknesses", cnt.get("weaknesses", 0)],
                    ["Opportunities", cnt.get("opportunities", 0)], ["Threats", cnt.get("threats", 0)],
                    ["Watch list", cnt.get("watch_list", 0)]]
            _table(s, ["SWOT quadrant", "Items"], rows, 6.2, 2.9, 5.2, col_w=[3.4, 1.8])
    else:
        _tb(s, 0.6, 2.4, 11.5, 1.2,
            "No closed assessment cycle yet — the Composite Excellence Index and SWOT "
            "appear here once the organization completes its first assessment.",
            size=15, color=MUTED, italic=True)

    # 10. Recommendations + dispositions
    s = _blank(prs); _header(s, "Recommendations", "AXIOM levers & decisions", logo_png)
    rows = []
    for r in (extras.get("recommendations") or [])[:8]:
        disp = r.get("disposition", "none")
        ini = r.get("initiative") or {}
        d = {"none": "—", "adopted": f"Adopted → {ini.get('ref', '')}"
             + (f", {ini.get('status')}" if ini.get("status") else "")
             + (f", {_rag_word(ini.get('rag'))}" if ini.get("rag") else ""),
             "parked": f"Parked → {ini.get('ref', '')}", "dismissed": "Dismissed"}.get(disp, disp)
        rows.append([r.get("title", "")[:48], f"{r.get('expected_ev_impact', 0):+,.1f}", d])
    if rows:
        _table(s, ["Recommendation", "EV impact", "Disposition"], rows, 0.6, 2.0, 12.1,
               col_w=[6.6, 2.0, 3.5])
    else:
        _tb(s, 0.6, 2.3, 11, 0.6, "No value-creating recommendations for the active dataset.",
            size=13, color=MUTED)

    # 11. Initiatives status board
    s = _blank(prs); _header(s, "Execution", "Key initiatives status board", logo_png)
    rows = [[i.get("ref_code"), (i.get("current_priority") or "").title(),
             _rag_word(i.get("rag")), i.get("owner_name") or "—",
             (i.get("status") or "").replace("_", " ").title()]
            for i in (extras.get("initiatives") or [])[:10]]
    if rows:
        _table(s, ["Ref", "Band", "RAG", "Owner", "Status"], rows, 0.6, 2.0, 12.1,
               col_w=[1.3, 2.0, 1.6, 4.0, 3.2])
    else:
        _tb(s, 0.6, 2.3, 11, 0.6, "No initiatives yet.", size=13, color=MUTED)

    # 12. Closing
    s = _blank(prs)
    _tb(s, 0.9, 2.6, 11.5, 0.9, "AXIOM Dynamics", size=32, color=BRASS, bold=True)
    _tb(s, 0.9, 3.6, 11.5, 1.4, [
        "This deck was generated from your active dataset by AXIOM's certified engines.",
        issued_line(issued, meta.get("dataset_version")),
        "support@axiomdynamics.app · axiomdynamics.app"], size=14, color=MUTED)

    out = io.BytesIO(); prs.save(out)
    return out.getvalue()


# ============================================================================
# Board Report PDF — the polished, multi-page report in report_pdf.py is the
# ONE AND ONLY board-report PDF builder. The minimal 7f dark builder that used
# to live here was removed (quarantined) so nothing can select it — one PDF
# truth. build_pdf simply delegates to it.
# ============================================================================
def build_pdf(report: dict, extras: dict, meta: dict) -> bytes:
    from .report_pdf import build_board_pdf
    return build_board_pdf(report, extras, meta)

# ============================================================================
# Comprehensive board presentation (7f revision) — mirrors the webapp nav
# ============================================================================
def _hist_series(data, block, line):
    try:
        hyears = [int(y) for y in data["periods"]["historical"]]
        blk = (data.get(block) or {}).get(line, {})
        return [{"year": y, "value": blk.get(str(y)), "kind": "hist"} for y in hyears
                if blk.get(str(y)) is not None]
    except Exception:
        return []


def _fcst_series(report, line):
    S = {s["id"]: s for s in report.get("sections", [])}
    out = []
    for st in (S.get("proforma", {}).get("statements") or []):
        v = ((st.get("stochastic") or {}).get(line) or {}).get("expected")
        if v is not None:
            out.append({"year": st["year"], "value": v, "kind": "fcst"})
    return out


def _margin_series(data, report):
    pts = []
    for p in _hist_series(data, "income_statement", "revenue"):
        eb = (data.get("income_statement") or {}).get("ebit", {}).get(str(p["year"]))
        if eb is not None and p["value"]:
            pts.append({"year": p["year"], "value": 100 * eb / p["value"], "kind": "hist"})
    S = {s["id"]: s for s in report.get("sections", [])}
    for st in (S.get("proforma", {}).get("statements") or []):
        r = ((st.get("stochastic") or {}).get("revenue") or {}).get("expected")
        e = ((st.get("stochastic") or {}).get("ebit") or {}).get("expected")
        if r and e is not None:
            pts.append({"year": st["year"], "value": 100 * e / r, "kind": "fcst"})
    return pts


class _CDeck:
    def __init__(self, prs, logo, company):
        self.prs, self.logo, self.company = prs, logo, company
        self.page, self.section = 0, "Overview"

    def _footer(self, s):
        self.page += 1
        _tb(s, 0.6, 7.08, 9.0, 0.3, f"AXIOM  ·  {self.company}  ·  Board Presentation",
            size=8, color=MUTED)
        _tb(s, 10.3, 7.08, 2.5, 0.3, f"{self.section}  ·  {self.page}", size=8, color=MUTED,
            align=PP_ALIGN.RIGHT)

    def content(self, kicker, title):
        s = _blank(self.prs); self._footer(s)
        _header(s, kicker, title, self.logo)
        return s

    def divider(self, name, subtitle=""):
        self.section = name
        s = _blank(self.prs); self.page += 1
        _rule(s, 0.9, 3.05, 3.4, color=BRASS, h=0.05)
        _tb(s, 0.9, 3.2, 11.5, 1.0, name, size=44, color=BRASS, bold=True)
        if subtitle:
            _tb(s, 0.94, 4.35, 11.5, 0.6, subtitle, size=16, color=MUTED)
        _tb(s, 10.3, 7.08, 2.5, 0.3, str(self.page), size=8, color=MUTED, align=PP_ALIGN.RIGHT)
        return s

    def placeholder(self, kicker, title, message):
        s = self.content(kicker, title)
        _panel(s, 1.1, 2.7, 11.1, 2.7)
        _tb(s, 1.5, 3.15, 10.3, 0.5, "Not available yet", size=13, color=BRASS, bold=True)
        _tb(s, 1.5, 3.65, 10.3, 1.6, message, size=15, color=MUTED, italic=True)
        return s


def _fmt2(v, pct=False, money=False):
    if v is None:
        return "—"
    if pct:
        return f"{v*100:.1f}%"
    return f"{v:,.1f}"


def build_pptx_comprehensive(report, extras, meta, data=None) -> bytes:
    extras = extras or {}; data = data or {}
    prs = Presentation()
    prs.slide_width = Emu(int(SW * EMU_IN)); prs.slide_height = Emu(int(SH * EMU_IN))
    S = {s["id"]: s for s in report.get("sections", [])}
    company = report.get("company", {})
    cname = meta.get("company_name", company.get("name", "Company"))
    cur = company.get("currency", "")
    issued = meta["issued_at"]
    logo = _logo_from_meta(meta)
    D = _CDeck(prs, logo, cname)

    # ---- Cover ----
    s = _blank(prs); D.page += 1
    if logo:
        sz = _px_size(logo); w, h = _fit_box(sz[0], sz[1], 5.0, 1.8) if sz else (4.2, 1.5)
        try:
            s.shapes.add_picture(io.BytesIO(logo), Inches(0.9), Inches(1.2), Inches(w), Inches(h))
        except Exception:
            logo = None; D.logo = None
    if not logo:
        _tb(s, 0.9, 1.9, 11.5, 0.5, "AXIOM", size=18, color=BRASS, bold=True)
    top = 3.4 if _logo_from_meta(meta) else 2.6
    _tb(s, 0.9, top, 11.5, 1.2, cname, size=40, color=IVORY, bold=True)
    _tb(s, 0.9, top + 1.3, 11.5, 0.6, "Comprehensive Board Presentation", size=22, color=MUTED)
    _rule(s, 0.94, top + 2.1, 4.0)
    _tb(s, 0.9, top + 2.3, 11.5, 0.5, issued_line(issued, meta.get("dataset_version")), size=13, color=MUTED)
    _tb(s, 0.9, 6.95, 11.5, 0.4, "Confidential — prepared by AXIOM Dynamics", size=10, color=MUTED, italic=True)

    # ---- Agenda ----
    s = D.content("Agenda", "What this deck covers")
    _tb(s, 0.7, 2.0, 5.9, 4.6, [
        "DIAGNOSE", "  • Dashboard & Reports", "  • SWOT Analysis",
        "  • Corporate Effectiveness (CEI)", "  • Valuation", "  • Financial Forecasts",
        "  • Risk Analysis", "  • Benchmarking"], size=14, color=IVORY)
    _tb(s, 6.9, 2.0, 5.9, 4.6, [
        "OPTIMIZE", "  • Enterprise Optimization", "  • Dynamics & Simulation",
        "  • Scenario Analysis", "  • Target-State Planner", "  • Key Initiatives",
        "  • Discussion", "  • Twin Monitoring"], size=14, color=IVORY)

    # ---- How to read ----
    s = D.content("How to read this deck", "Bands, RAG, and valuation lenses")
    for i, (t, b) in enumerate([
        ("P05–P95 bands", "Charts show a median with a shaded uncertainty band — the range within which the outcome falls 90% of the time."),
        ("RAG status", "Green = on track / strong; Amber = watch; Red = off track / weak. Applied to scores, CSFs, and initiatives."),
        ("Three valuation lenses", "Intrinsic DCF, a Monte-Carlo distribution, and comparable multiples — read together, never in isolation."),
        ("Dispositions", "AXIOM recommendations carry a decision chip: Adopted → an initiative ref, Parked, or Dismissed.")]):
        cy = 2.0 + i * 1.18
        _panel(s, 0.7, cy, 11.9, 1.02)
        _tb(s, 0.95, cy + 0.12, 3.3, 0.8, t, size=13, color=BRASS, bold=True)
        _tb(s, 4.3, cy + 0.12, 8.1, 0.8, b, size=11.5, color=IVORY)

    # ================= DIAGNOSE =================
    D.divider("Diagnose", "Where the company stands today")

    # --- Dashboard & Reports ---
    D.section = "Dashboard"
    sm = S.get("summary", {}); sc = sm.get("scorecard", {}); di = S.get("diagnostic", {})
    kpis = di.get("kpi_strip") or []
    s = D.content("Dashboard & Reports", "Key performance indicators")
    rows = [[k.get("kpi"), _fmt2(k.get("current")), _fmt2(k.get("previous")), _fmt2(k.get("trend"), pct=True)]
            for k in kpis[:9]]
    if rows:
        _table(s, ["KPI", "Current", "Prior", "Trend"], rows, 0.7, 2.0, 8.6, col_w=[3.4, 1.7, 1.7, 1.8])
    _tb(s, 9.5, 2.0, 3.3, 4.6, report.get("units_note", "Figures in $ millions."), size=10, color=MUTED, italic=True)

    s = D.content("Dashboard & Reports", "Enterprise health index")
    hi = (di.get("health") or {}).get("index", sc.get("health_index"))
    _tb(s, 0.7, 2.0, 6.0, 1.0, f"{_fmt2(hi)} / 100", size=40, color=BRASS, bold=True)
    comps = (di.get("health") or {}).get("components") or {}
    if isinstance(comps, dict) and comps:
        _pic(s, chart_bars(list(comps.keys()), [comps[k] for k in comps],
                           "Health components", horizontal=True, fmt="{:.1f}"), 0.7, 3.0, 7.6, 3.6)
    _tb(s, 8.6, 2.0, 4.2, 4.4, [f"Risk grade: {sc.get('risk_grade','—')}",
        f"Optimization: {sc.get('optimization_status','—')}",
        f"Distance to optimum: {_fmt2((di.get('health') or {}).get('distance_to_optimum'))}"], size=12, color=IVORY)

    s = D.content("Dashboard & Reports", "Enterprise value summary")
    hm = sm.get("headline_metric", {})
    _tb(s, 0.7, 2.0, 6.0, 0.4, hm.get("label", "Enterprise Value"), size=13, color=MUTED)
    _tb(s, 0.7, 2.4, 6.0, 1.0, f"{hm.get('value', 0):,.0f} {cur}", size=38, color=BRASS, bold=True)
    va = S.get("valuation", {})
    _pic(s, chart_valuation_lenses(va.get("dcf"), va.get("real_options")), 0.7, 3.5, 7.4, 3.2)
    _tb(s, 8.4, 2.2, 4.4, 4.2, [f"• {w}" for w in sm.get("four_answers", [])[:4]], size=11.5, color=IVORY)

    s = D.content("Dashboard & Reports", "Revenue trajectory")
    _pic(s, chart_trend(_hist_series(data, "income_statement", "revenue") + _fcst_series(report, "revenue"),
                        "Revenue — actual & forecast ($M)"), 0.7, 2.0, 11.9, 4.6)

    s = D.content("Dashboard & Reports", "EBITDA trajectory")
    _pic(s, chart_trend(_hist_series(data, "income_statement", "ebitda") + _fcst_series(report, "ebitda"),
                        "EBITDA — actual & forecast ($M)"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Dashboard & Reports", "Free cash flow")
    _pic(s, chart_trend(_fcst_series(report, "fcff"), "Free cash flow to firm — forecast ($M)"), 0.7, 2.0, 11.9, 4.6)

    s = D.content("Dashboard & Reports", "Margin trajectory & alerts")
    _pic(s, chart_trend(_margin_series(data, report), "EBIT margin (%)"), 0.6, 2.0, 6.6, 4.4)
    runway = S.get("outlook", {}).get("cash_runway", {})
    alerts = []
    if runway.get("burning_cash"):
        alerts.append("• Cash is burning — monitor runway")
    if (runway.get("p_cash_below_zero_ever") or 0) > 0.05:
        alerts.append(f"• P(cash<0) ever: {_fmt2(runway.get('p_cash_below_zero_ever'), pct=True)}")
    for hm2 in (S.get("appendix", {}).get("risk_heat_map") or []):
        if hm2.get("rag") == "red":
            alerts.append(f"• {hm2.get('category')} risk elevated")
    _tb(s, 7.5, 2.1, 5.3, 0.4, "Alerts", size=14, color=BRASS, bold=True)
    _tb(s, 7.5, 2.55, 5.3, 4.0, alerts or ["• No material alerts"], size=12, color=IVORY)

    # --- SWOT ---
    D.section = "SWOT"
    swot = extras.get("swot")
    if swot and swot.get("has_data"):
        cnt = swot.get("counts", {})
        s = D.content("SWOT Analysis", "Quadrant overview")
        quads = [("Strengths", "strengths", GREEN), ("Weaknesses", "weaknesses", RED),
                 ("Opportunities", "opportunities", BRASS), ("Threats", "threats", AMBER)]
        for i, (lab, key, col) in enumerate(quads):
            cx = 0.7 + (i % 2) * 6.1; cy = 2.0 + (i // 2) * 2.35
            _panel(s, cx, cy, 5.8, 2.15)
            _tb(s, cx + 0.25, cy + 0.15, 5.3, 0.4, f"{lab}  ({cnt.get(key, 0)})", size=14, color=col, bold=True)
            items = swot.get(key, [])[:4]
            _tb(s, cx + 0.25, cy + 0.6, 5.3, 1.5,
                [f"• {e.get('item_code')} {e.get('title','')[:34]}" for e in items] or ["—"],
                size=10.5, color=IVORY)
        for lab, key, col in quads:
            items = swot.get(key, [])
            if not items:
                continue
            s = D.content("SWOT Analysis", lab)
            rows = [[e.get("item_code"), (e.get("title") or "")[:40], _fmt2(e.get("mean")),
                     (e.get("score_rag") or "—").title(), (e.get("theme") or "—")[:34],
                     ", ".join(li.get("ref_code", "") for li in (e.get("linked_initiatives") or [])) or "—"]
                    for e in items[:8]]
            _table(s, ["Item", "Title", "Score", "RAG", "Theme", "Initiatives"], rows, 0.6, 2.0, 12.2,
                   col_w=[1.1, 3.6, 1.2, 1.2, 3.3, 1.8], fsize=10)
        s = D.content("SWOT Analysis", "Watch list")
        wl = swot.get("watch_list", [])
        rows = [[e.get("item_code"), (e.get("title") or "")[:48], _fmt2(e.get("mean")), (e.get("score_rag") or "—").title()] for e in wl[:10]]
        if rows:
            _table(s, ["Item", "Title", "Score", "RAG"], rows, 0.7, 2.0, 11.6, col_w=[1.3, 6.0, 1.8, 2.0])
        else:
            _tb(s, 0.7, 2.2, 11, 0.5, "No mid-band watch-list items.", size=13, color=MUTED)
    else:
        D.placeholder("SWOT Analysis", "SWOT Analysis",
                      "SWOT is derived from your latest closed assessment cycle. It appears here once the "
                      "organization completes its first assessment — strengths, weaknesses, opportunities and "
                      "threats, each with score, RAG, sentiment theme and linked initiatives.")

    # --- Corporate Effectiveness (CEI) ---
    D.section = "Corporate Effectiveness"
    cei = extras.get("cei")
    if cei and cei.get("cei") is not None:
        s = D.content("Corporate Effectiveness", "Composite Excellence Index")
        _tb(s, 0.7, 2.1, 6, 1.0, f"CEI  {cei['cei']:.1f} / 10", size=40, color=BRASS, bold=True)
        trend = cei.get("trend") or []
        if len(trend) >= 2:
            _pic(s, chart_trend([{"year": i + 1, "value": t.get("cei"), "kind": "hist"} for i, t in enumerate(trend)],
                                "CEI trend (by cycle)"), 0.7, 3.1, 7.4, 3.5)
        _tb(s, 8.4, 2.1, 4.4, 4.3, f"{cei.get('n_participants', 0)} participants", size=13, color=IVORY)
        s = D.content("Corporate Effectiveness", "13-axis excellence radar")
        _pic(s, chart_radar(cei.get("l1_subscores")), 3.6, 1.75, 6.1, 5.2)
        subs = cei.get("l1_subscores") or []
        for half, (lo, hi) in enumerate([(0, 7), (7, 14)]):
            s = D.content("Corporate Effectiveness", f"Subscores ({half+1}/2)")
            rows = [[o.get("title", "")[:40], _fmt2(o.get("score")), _fmt2((o.get("dispersion") or {}).get("std"))]
                    for o in subs[lo:hi]]
            if rows:
                _table(s, ["Category", "Score", "Dispersion"], rows, 0.7, 2.0, 11.6, col_w=[7.0, 2.0, 2.6])
        s = D.content("Corporate Effectiveness", "Consensus & comment themes")
        disp = cei.get("item_dispersion") or {}
        hi_disp = sorted(disp.items(), key=lambda kv: -(kv[1].get("std") or 0))[:6]
        _tb(s, 0.7, 2.0, 12, 0.4, "Widest dispersion (least consensus):", size=13, color=BRASS, bold=True)
        _tb(s, 0.7, 2.5, 12, 3.0, [f"• {v.get('title', k)}: σ={_fmt2(v.get('std'))}" for k, v in hi_disp] or ["—"],
            size=12, color=IVORY)
    else:
        D.placeholder("Corporate Effectiveness", "Corporate Effectiveness",
                      "The Composite Excellence Index scores your organization across 13 capability domains with "
                      "a 361-point taxonomy. It becomes available after your first assessment cycle closes — with a "
                      "trend, a 13-axis radar, subscores, a consensus map and comment-theme sentiment.")

    # --- Valuation ---
    D.section = "Valuation"
    dcf = va.get("dcf", {}); mc = dcf.get("monte_carlo", {})
    s = D.content("Valuation", "Three independent lenses")
    _pic(s, chart_valuation_lenses(dcf, va.get("real_options")), 0.8, 1.95, 8.0, 4.6)
    _tb(s, 9.1, 2.0, 3.6, 4.5, [f"DCF: {dcf.get('enterprise_value', 0):,.0f}",
        f"MC mean: {mc.get('mean', 0):,.0f}", f"MC CVaR95: {mc.get('cvar95', 0):,.0f}",
        f"WACC: {_fmt2(dcf.get('wacc'), pct=True)}"], size=13, color=IVORY)
    s = D.content("Valuation", "Discounted cash flow — detail")
    br = dcf.get("bridge") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in list(br.items())[:8]] if isinstance(br, dict) else []
    if rows:
        _table(s, ["Bridge component", "Value ($M)"], rows, 0.7, 2.0, 8.0, col_w=[5.4, 2.6])
    _tb(s, 9.0, 2.0, 3.8, 4.4, [f"Enterprise value: {dcf.get('enterprise_value', 0):,.0f}",
        f"Equity value: {dcf.get('equity_value', 0):,.0f}",
        f"Post-DLOM: {_fmt2(dcf.get('equity_value_post_dlom'))}"], size=12, color=IVORY)
    s = D.content("Valuation", "Monte Carlo distribution")
    _pic(s, chart_mc_hist(mc.get("histogram"), mc.get("percentiles")), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Valuation", "Value drivers")
    _pic(s, chart_tornado(S.get("actions", {}).get("recommendations")), 0.7, 1.95, 11.9, 4.8)
    s = D.content("Valuation", "Real-options premium")
    ro = va.get("real_options", {}); opts = ro.get("options", {})
    rows = [[k.title(), _fmt2((opts.get(k) or {}).get("flexibility_value"))] for k in opts]
    if rows:
        _table(s, ["Option", "Flexibility value ($M)"], rows, 0.7, 2.0, 8.0, col_w=[4.0, 4.0])
    _tb(s, 9.0, 2.0, 3.8, 3.5, f"Total flexibility:\n{_fmt2(ro.get('total_flexibility_value'))} {cur}", size=14, color=BRASS, bold=True)
    s = D.content("Valuation", "Interest-rate sensitivity")
    rs = va.get("rate_sensitivity") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in rs.items() if isinstance(v, (int, float))][:8]
    if rows:
        _table(s, ["Measure", "Value"], rows, 0.7, 2.0, 8.5, col_w=[5.3, 3.2])
    else:
        _tb(s, 0.7, 2.2, 11, 0.6, "Duration/convexity of value to the discount rate — see the appendix.", size=13, color=MUTED)

    # --- Financial Forecasts ---
    D.section = "Financial Forecasts"
    pf = S.get("proforma", {}); stmts = pf.get("statements") or []
    def _band_table(line_keys, title):
        s = D.content("Financial Forecasts", title)
        headers = ["Line"] + [str(st["year"]) for st in stmts]
        rows = []
        for lk, label in line_keys:
            cells = [label]
            for st in stmts:
                d = (st.get("stochastic") or {}).get(lk) or {}
                cells.append(f"{_fmt2(d.get('p05'))}–{_fmt2(d.get('p95'))}" if d else "—")
            rows.append(cells)
        if rows:
            _table(s, headers, rows, 0.5, 2.0, 12.3, fsize=9)
        _tb(s, 0.5, 6.4, 12, 0.4, "Ranges are P05–P95 ($M).", size=9, color=MUTED, italic=True)
    _band_table([("revenue", "Revenue"), ("ebit", "EBIT"), ("ebitda", "EBITDA"), ("net_income", "Net income")],
                "Pro-forma income statement (P05–P95)")
    _band_table([("total_assets", "Total assets"), ("equity", "Equity"), ("cash", "Cash")],
                "Pro-forma balance sheet (P05–P95)")
    _band_table([("cfo", "Cash from ops"), ("fcff", "FCFF"), ("fcfe", "FCFE")],
                "Pro-forma cash flow (P05–P95)")
    sb = S.get("outlook", {}).get("simulation_baseline", {})
    s = D.content("Financial Forecasts", "Revenue fan"); _pic(s, chart_fan(sb.get("revenue_fan"), "Revenue — P05–P95"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Financial Forecasts", "Free-cash-flow fan"); _pic(s, chart_fan(sb.get("fcff_fan"), "FCFF — P05–P95"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Financial Forecasts", "Cash fan"); _pic(s, chart_fan(sb.get("cash_fan"), "Cash — P05–P95"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Financial Forecasts", "Comprehensive income")
    _tb(s, 0.7, 2.0, 12, 0.5, f"Accounting framework: {pf.get('accounting_framework', '—')}", size=13, color=BRASS, bold=True)
    _tb(s, 0.7, 2.6, 12, 3.8, [f"• {w}" for w in (pf.get("narrative") or [])[:4]], size=12.5, color=IVORY)
    s = D.content("Financial Forecasts", "Plan attainment")
    pa = S.get("outlook", {}).get("plan_attainment") or {}
    rows = [["Revenue target (yr 1)", _fmt2(pa.get("p_revenue_target"), pct=True)],
            ["EBIT-margin target (yr 1)", _fmt2(pa.get("p_margin_target"), pct=True)],
            ["FCFF target (yr 1)", _fmt2(pa.get("p_fcff_target"), pct=True)],
            ["All-three (yr 1)", _fmt2(pa.get("p_all_three") or pa.get("p_all_thresholds"), pct=True)]]
    _table(s, ["Plan target", "P(meet or beat)"], rows, 0.7, 2.0, 9.0, col_w=[5.4, 3.6])

    # --- Risk Analysis ---
    D.section = "Risk Analysis"
    rg = di.get("risk_grade", {}); cov = S.get("outlook", {}).get("coverage", {})
    s = D.content("Risk Analysis", "Risk grade & indicators")
    _tb(s, 0.7, 2.0, 4, 1.0, f"Grade  {rg.get('grade', sc.get('risk_grade','—'))}", size=36, color=BRASS, bold=True)
    rows = [["Total debt", _fmt2(cov.get("total_debt"))],
            ["Distance to default (σ)", _fmt2(cov.get("distance_to_default_sigmas"))],
            ["P(EV < debt)", _fmt2(cov.get("p_ev_below_debt"), pct=True)],
            ["P(cash<0) baseline", _fmt2(cov.get("p_cash_below_zero_baseline"), pct=True)],
            ["P(cash<0) recession", _fmt2(cov.get("p_cash_below_zero_recession"), pct=True)]]
    _table(s, ["Indicator", "Value"], rows, 0.7, 3.1, 7.6, col_w=[4.6, 3.0])
    s = D.content("Risk Analysis", "Distance to default")
    _tb(s, 0.7, 2.2, 11.8, 1.2, f"{_fmt2(cov.get('distance_to_default_sigmas'))} σ to the default barrier",
        size=26, color=IVORY, bold=True)
    _tb(s, 0.7, 3.6, 11.8, 2.4, cov.get("method", ""), size=12, color=MUTED)
    s = D.content("Risk Analysis", "Outcome distribution (stress)")
    _pic(s, chart_mc_hist(mc.get("histogram"), mc.get("percentiles")), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Risk Analysis", "Risk drivers")
    heat = S.get("appendix", {}).get("risk_heat_map") or []
    if heat:
        _pic(s, chart_bars([h.get("category") for h in heat], [h.get("score") for h in heat],
                           "Risk heat map (variance share)",
                           colors=[{"red": RED, "amber": AMBER, "green": GREEN}.get(h.get("rag"), BRASS) for h in heat],
                           horizontal=True, fmt="{:.0f}"), 0.7, 2.0, 11.6, 4.5)
    s = D.content("Risk Analysis", "Covenant headroom")
    covs = S.get("appendix", {}).get("covenants") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in (covs.items() if isinstance(covs, dict) else []) if isinstance(v, (int, float))][:8]
    if rows:
        _table(s, ["Covenant", "Value"], rows, 0.7, 2.0, 9.0, col_w=[5.4, 3.6])
    else:
        _tb(s, 0.7, 2.2, 11, 0.6, "No debt covenants configured for this company.", size=13, color=MUTED)
    s = D.content("Risk Analysis", "Extreme-value tail")
    evt = S.get("appendix", {}).get("extreme_value_tail") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in (evt.items() if isinstance(evt, dict) else []) if isinstance(v, (int, float))][:8]
    if rows:
        _table(s, ["Tail measure", "Value"], rows, 0.7, 2.0, 9.0, col_w=[5.4, 3.6])
    else:
        _tb(s, 0.7, 2.2, 11, 0.6, "Extreme-value (peaks-over-threshold) tail statistics — see appendix.", size=13, color=MUTED)

    # --- Benchmarking ---
    D.section = "Benchmarking"
    bench = di.get("benchmark", {}); bkpis = extras.get("benchmark_kpis") or []
    if bench and (bench.get("index") is not None or bkpis):
        s = D.content("Benchmarking", "Versus industry")
        _tb(s, 0.7, 2.0, 5, 1.0, f"BPI  {_fmt2(bench.get('index'))}", size=36, color=BRASS, bold=True)
        _tb(s, 0.7, 3.1, 12, 2.6, (bench.get("narrative") or "")[:400], size=13, color=IVORY)
        grp = [k for k in bkpis if k.get("score") is not None][:9]
        if grp:
            s = D.content("Benchmarking", "Scores vs sector")
            _pic(s, chart_bars([k.get("label", k.get("kpi", ""))[:16] for k in grp],
                               [k.get("score") for k in grp], "Benchmark scores (1.0 = in line with sector)",
                               colors=[{"red": RED, "amber": AMBER, "green": GREEN}.get(k.get("rag"), BRASS) for k in grp],
                               horizontal=True, fmt="{:.2f}"), 0.7, 2.0, 11.7, 4.6)
            s = D.content("Benchmarking", "Benchmark detail")
            def _bv(k, key):
                v = k.get(key)
                return _fmt2(v, pct=True) if k.get("format") == "percent" else _fmt2(v)
            rows = [[k.get("label", k.get("kpi")), _bv(k, "actual"), _bv(k, "benchmark"), (k.get("rag") or "—").title()] for k in grp]
            _table(s, ["KPI", "Company", "Sector", "RAG"], rows, 0.7, 2.0, 11.0, col_w=[4.0, 2.3, 2.3, 2.4], fsize=10)
    else:
        D.placeholder("Benchmarking", "Benchmarking",
                      "Peer benchmarking compares your margins, growth, ROIC and leverage against a curated sector "
                      "set. It appears here once a sector peer set is selected for your company.")

    # ================= OPTIMIZE =================
    D.divider("Optimize", "What should change, and by how much")

    # --- Enterprise Optimization ---
    D.section = "Enterprise Optimization"
    ac = S.get("actions", {}); bd = S.get("best_decision", {})
    s = D.content("Enterprise Optimization", "Optimal levers")
    plan = ac.get("optimizer_plan") or []
    rows = [[f"Step {p.get('step')}", _fmt2(p.get("growth"), pct=True), _fmt2(p.get("net_borrowing_pct_rev"), pct=True),
             _fmt2(p.get("revenue_target"))] for p in plan[:8]]
    if rows:
        _table(s, ["Step", "Growth", "Net borrowing %", "Revenue target"], rows, 0.7, 2.0, 11.0, col_w=[2.0, 2.6, 3.0, 3.4])
    _tb(s, 0.7, 6.2, 12, 0.5, f"Optimization uplift: {_fmt2(ac.get('optimization_uplift'))} {cur}", size=13, color=BRASS, bold=True)
    s = D.content("Enterprise Optimization", "Upside bridge")
    ud = ac.get("uplift_derivation") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in (ud.items() if isinstance(ud, dict) else []) if isinstance(v, (int, float))][:8]
    if rows:
        _table(s, ["Bridge step", "Value ($M)"], rows, 0.7, 2.0, 9.0, col_w=[5.4, 3.6])
    _tb(s, 0.7, 6.2, 12, 0.5, f"Total optimization uplift: {_fmt2(ac.get('optimization_uplift'))} {cur}", size=13, color=BRASS, bold=True)
    s = D.content("Enterprise Optimization", "Value–risk frontier")
    _pic(s, chart_frontier((bd.get("frontier") or {}).get("points"), (bd.get("frontier") or {}).get("recommended")), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Enterprise Optimization", "Optimization ladder")
    _pic(s, chart_trend([{"year": p.get("step"), "value": p.get("revenue_target"), "kind": "fcst"} for p in plan],
                        "Revenue target by optimization step"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Enterprise Optimization", "Ranked recommendations")
    rr = []
    for r in (extras.get("recommendations") or [])[:8]:
        ini = r.get("initiative") or {}
        chip = {"adopted": f"Adopted → {ini.get('ref','')}" + (f", {ini.get('status')}" if ini.get('status') else "") + (f", {(ini.get('rag') or '').title()}" if ini.get('rag') else ""),
                "parked": f"Parked → {ini.get('ref','')}", "dismissed": "Dismissed", "none": "—"}.get(r.get("disposition"), r.get("disposition"))
        rr.append([(r.get("title") or "")[:46], f"{r.get('expected_ev_impact', 0):+,.1f}", chip])
    if rr:
        _table(s, ["Recommendation", "EV impact", "Disposition"], rr, 0.6, 2.0, 12.2, col_w=[6.6, 2.0, 3.6])
    s = D.content("Enterprise Optimization", "Shadow prices")
    sp = bd.get("shadow_prices") or {}
    rows = [[k.replace("_", " ").title(), _fmt2(v)] for k, v in sp.items()]
    if rows:
        _table(s, ["Constraint", "Shadow price ($M)"], rows, 0.7, 2.0, 9.0, col_w=[5.6, 3.4])

    # --- Dynamics & Simulation ---
    D.section = "Dynamics & Simulation"
    yrs = sb.get("years") or []
    s = D.content("Dynamics & Simulation", "Simulated revenue paths")
    _pic(s, chart_paths((sb.get("sample_paths") or {}).get("revenue"), yrs, "Simulated revenue paths ($M)"), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Dynamics & Simulation", "Simulated cash paths")
    _pic(s, chart_paths((sb.get("sample_paths") or {}).get("cash"), yrs, "Simulated cash paths ($M)", color=GREEN), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Dynamics & Simulation", "Simulated FCFF paths")
    _pic(s, chart_paths((sb.get("sample_paths") or {}).get("fcff"), yrs, "Simulated FCFF paths ($M)", color=AMBER), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Dynamics & Simulation", "Outcome distribution")
    _pic(s, chart_mc_hist(mc.get("histogram"), mc.get("percentiles")), 0.7, 2.0, 11.9, 4.6)

    # --- Scenario Analysis ---
    D.section = "Scenario Analysis"
    rec_sim = S.get("outlook", {}).get("simulation_recession", {})
    s = D.content("Scenario Analysis", "Baseline vs recession")
    _pic(s, chart_fan(sb.get("revenue_fan"), "Baseline revenue — P05–P95"), 0.6, 2.0, 6.2, 4.4)
    _pic(s, chart_fan(rec_sim.get("revenue_fan"), "Recession revenue — P05–P95"), 6.9, 2.0, 6.2, 4.4)
    s = D.content("Scenario Analysis", "Per-lever swings")
    _pic(s, chart_frontier((bd.get("frontier") or {}).get("points"), (bd.get("frontier") or {}).get("recommended")), 0.7, 2.0, 11.9, 4.6)
    s = D.content("Scenario Analysis", "Scenario summary")
    _tb(s, 0.7, 2.1, 12, 3.5, [
        f"• Recession P(cash<0) ever: {_fmt2(rec_sim.get('p_cash_below_zero'), pct=True)}",
        "• The frontier traces value against tail-safety across capital structures.",
        "• The interactive What-If Studio in AXIOM lets you move any lever live."], size=13, color=IVORY)

    # --- Target-State Planner ---
    D.section = "Target-State Planner"
    if plan:
        s = D.content("Target-State Planner", "Current vs target")
        base_rev = _hist_series(data, "income_statement", "revenue")
        cur_rev = base_rev[-1]["value"] if base_rev else None
        tgt_rev = plan[-1].get("revenue_target")
        rows = [["Revenue", _fmt2(cur_rev), _fmt2(tgt_rev),
                 _fmt2((tgt_rev - cur_rev) if (cur_rev and tgt_rev) else None)],
                ["Optimization uplift", "—", _fmt2(ac.get("optimization_uplift")), _fmt2(ac.get("optimization_uplift"))]]
        _table(s, ["Metric", "Current", "Target", "Gap"], rows, 0.7, 2.0, 11.0, col_w=[3.6, 2.4, 2.4, 2.6])
        s = D.content("Target-State Planner", "Path to target")
        _pic(s, chart_trend([{"year": p.get("step"), "value": p.get("revenue_target"), "kind": "fcst"} for p in plan],
                            "Revenue target ladder"), 0.7, 2.0, 11.9, 4.6)
        s = D.content("Target-State Planner", "Gap closure")
        _tb(s, 0.7, 2.1, 12, 3.0, [
            f"• Target revenue: {_fmt2(tgt_rev)} {cur}",
            f"• Optimizer-available uplift: {_fmt2(ac.get('optimization_uplift'))} {cur}",
            "• Set explicit targets in the Target-State Planner to quantify each gap."], size=13, color=IVORY)
    else:
        D.placeholder("Target-State Planner", "Target-State Planner",
                      "Define a desired future state (revenue, margin, capital structure) and AXIOM quantifies each "
                      "gap and maps it to the value-creating lever that closes it.")

    # --- Key Initiatives ---
    D.section = "Key Initiatives"
    inits = extras.get("initiatives") or []
    if inits:
        s = D.content("Key Initiatives", "Register board")
        rows = [[i.get("ref_code"), (i.get("current_priority") or "").title(),
                 (i.get("rag") or "—").title(), i.get("owner_name") or "—",
                 (i.get("status") or "").replace("_", " ").title()] for i in inits[:12]]
        _table(s, ["Ref", "Band", "RAG", "Owner", "Status"], rows, 0.6, 2.0, 12.2, col_w=[1.3, 2.0, 1.6, 4.0, 3.3])
        csf = extras.get("csf_health") or {}
        s = D.content("Key Initiatives", "CSF health summary")
        rows = [["Holding", csf.get("holding", 0)], ["At risk", csf.get("at_risk", 0)], ["Broken", csf.get("broken", 0)]]
        _table(s, ["CSF status", "Count"], rows, 0.7, 2.0, 7.0, col_w=[4.0, 3.0])
        s = D.content("Key Initiatives", "Expected vs realized")
        rows = [[i.get("ref_code"), _fmt2(i.get("expected_impact_amount")), _fmt2(i.get("actual_impact_amount"))]
                for i in inits[:12] if i.get("expected_impact_amount") is not None or i.get("actual_impact_amount") is not None]
        if rows:
            _table(s, ["Ref", "Expected ($M)", "Realized ($M)"], rows, 0.7, 2.0, 9.0, col_w=[2.0, 3.5, 3.5])
        else:
            _tb(s, 0.7, 2.2, 11, 0.5, "No settlements recorded yet.", size=13, color=MUTED)
    else:
        D.placeholder("Key Initiatives", "Key Initiatives",
                      "The execution register turns decisions into tracked initiatives — banded by priority, with a "
                      "RAG, an owner, critical success factors and expected-vs-realized impact. It populates as you "
                      "adopt recommendations and create initiatives.")

    # --- Discussion ---
    D.section = "Discussion"
    disc = extras.get("discussion") or {}
    if disc.get("threads") or disc.get("pending_proposals"):
        s = D.content("Discussion", "Activity & proposals")
        _tb(s, 0.7, 2.0, 12, 0.5, f"{disc.get('threads', 0)} threads · {disc.get('posts', 0)} posts · "
            f"{disc.get('pending_proposals', 0)} pending proposals", size=15, color=IVORY, bold=True)
        titles = disc.get("proposal_titles") or []
        _tb(s, 0.7, 2.8, 12, 3.5, [f"• {t}" for t in titles[:8]] or ["No pending proposals."], size=12, color=IVORY)
    else:
        D.placeholder("Discussion", "Discussion",
                      "The discussion forum threads every report, initiative and topic. Flagged posts become "
                      "proposals the board can adopt. It appears here once conversations begin.")

    # --- Twin Monitoring ---
    D.section = "Twin Monitoring"
    D.placeholder("Twin Monitoring", "Twin Monitoring",
                  "As actuals arrive, AXIOM compares them to the digital twin's projection — surfacing drift and "
                  "auto-recalibrating. This view activates once an actuals sync is attached to this company.")

    # ================= CLOSE =================
    D.section = "Appendix"
    s = D.content("Methodology", "How these numbers are produced")
    _tb(s, 0.7, 2.0, 12.0, 4.6, (S.get("appendix", {}).get("methodology") or
        "AXIOM runs a certified valuation, risk, optimization and simulation suite on your dataset. Every figure "
        "is reproducible from the active dataset version shown on the cover.")[:900], size=12.5, color=IVORY)
    s = D.content("Glossary", "Key terms")
    _tb(s, 0.7, 2.0, 12.0, 4.6, [
        "EV — Enterprise Value.   DCF — Discounted Cash Flow.   WACC — Weighted Average Cost of Capital.",
        "CVaR — Conditional Value at Risk (expected loss in the worst 5%).",
        "CEI — Composite Excellence Index.   RAG — Red / Amber / Green status.",
        "P05–P95 — the 5th to 95th percentile band of the simulated distribution.",
        "BPI — Benchmark Performance Index (100 = in line with sector)."], size=13, color=IVORY)
    s = _blank(prs); D.page += 1
    _tb(s, 0.9, 2.7, 11.5, 0.9, "AXIOM Dynamics", size=32, color=BRASS, bold=True)
    _tb(s, 0.9, 3.7, 11.5, 1.6, [
        "Generated from your active dataset by AXIOM's certified engines.",
        issued_line(issued, meta.get("dataset_version")),
        "support@axiomdynamics.app · axiomdynamics.app"], size=14, color=MUTED)

    out = io.BytesIO(); prs.save(out)
    return out.getvalue()
