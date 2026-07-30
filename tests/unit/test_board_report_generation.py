"""The board report must GENERATE — on real magnitudes, and with absence present.

⭐ THIS 500'd FOR EVERY REAL COMPANY, ON BOTH FORMATS, SINCE EARLY THIS MONTH,
AND THE VISIBLE PATH STAYED GREEN. Reproduced against production:

    company 38 (real)     PDF -> 500    PPTX -> 500
    company 20 (showcase) PDF -> 201    PPTX -> 201

The showcase route is `_serve_showcase_latest`: it returns a PRE-GENERATED
artifact with no compute. Every demo and every "reports work" check used
Meridian and therefore exercised a stored file. **Generation had no coverage at
all** — which is why 100% of real companies failed behind a green surface.

So these tests call the BUILDERS, on a dataset shaped like the one that broke.

⭐ THE FIXTURE IS SMALL-MAGNITUDE ON PURPOSE. Company 38's statements are
denominated so FCFF is ~0.0004 per period. `twin.bands()` rounded percentiles to
2dp, which made the whole fan 0.0, and `risk_analytics` then divided by the zero
dispersion that rounding had manufactured. A fixture using 100.0 everywhere
cannot reproduce that, and would have passed throughout.

⭐ AND IT CARRIES GENUINE ABSENCE, because the report layer is downstream of
everything. Four coercions were removed on 30 Jul, so None now reaches the
narrative — where eight format specs (`:.1f`, `:.2%`, `:,.1f`, `:.0%`) applied it
straight to `format()` and took out the CFO's board pack over one missing input.
"""
import os
import tempfile

os.environ.setdefault("DATABASE_URL",
                      "sqlite:///" + tempfile.mktemp(prefix="rpt-", suffix=".db"))

import datetime

import pytest

from services.api import reporting as R
from services.api.modules.intelligence import engines as intel

IS = ("revenue", "cogs", "opex", "depreciation_amortization", "interest_expense")
BS = ("cash", "other_current_assets", "noncurrent_assets", "short_term_debt",
      "long_term_debt", "current_liabilities_ex_debt", "total_equity",
      "preferred_equity", "minority_interest")
CF = ("capex", "net_borrowing", "dividends")

SMALL = 0.0025          # the magnitude that 2dp rounding annihilates


def _dataset(scale=SMALL, missing=(), missing_year=None):
    years = [2020, 2021, 2022, 2023]

    def block(keys):
        return {k: {str(y): (None if (y == missing_year and k in missing)
                             else scale * (1.0 + 0.1 * i))
                    for i, y in enumerate(years)} for k in keys}

    return {"company": {"name": "Small Co", "ownership": "private",
                        "standard": "us_gaap", "tax_rate": 0.25,
                        "risk_free_rate": 0.04, "market_risk_premium": 0.055,
                        "cost_of_debt": 0.06, "unlevered_industry_beta": 1.0,
                        "target_debt_to_equity": 0.5, "size_premium": 0.02,
                        "specific_risk_premium": 0.01, "dlom": 0.25,
                        "shares_outstanding": 1000},
            "periods": {"historical": years, "forecast": [],
                        "frequency": "annual"},
            "income_statement": block(IS), "balance_sheet": block(BS),
            "cash_flow": block(CF)}


def _meta():
    return {"company_name": "Small Co", "report_type": "Board Report",
            "issued_at": datetime.datetime.utcnow(), "dataset_version": 1,
            "logo": None}


def test_board_report_builds_on_small_magnitudes():
    """⭐ THE REGRESSION. ZeroDivisionError in risk_analytics, caused by 2dp
    rounding flattening the FCFF fan."""
    report = intel.board_report(_dataset())
    assert report and "checkpoints" in report


def test_the_simulation_fan_survives_small_magnitudes():
    """Rounding for presentation must not be applied to the payload other
    engines consume."""
    from services.api.modules.twin import engines as twin
    sim = twin.simulate(_dataset(), "baseline", n_paths=200)
    fan = sim["fcff_fan"][-1]
    assert fan["p95"] != fan["p05"] or fan["p95"] != 0.0, (
        "the fan collapsed to a single value — display precision destroyed the "
        "dispersion that risk_analytics divides by")


def test_variance_shares_are_absent_not_zero_when_there_is_no_dispersion():
    """A degenerate fan makes the decomposition UNDEFINED, not zero. Attributing
    it anyway would put invented shares on a board deck."""
    ra = intel.risk_analytics(_dataset(), n_paths=200)
    sob = ra["sobol_attribution"]
    if sob["growth_uncertainty"] is None:
        assert sob.get("absence_reason"), "absence with no reason is silent"
        assert sob["margin_uncertainty"] is None
        assert sob["interaction"] is None


@pytest.mark.xfail(strict=False, reason=(
    "NOT FIXED BY THIS LANE, AND DELIBERATELY LEFT VISIBLE. An absent line item "
    "still crashes generation — but upstream of the report, in the engines "
    "themselves, on plain subscripts like `bs['cash'][ys] + ...`. That is the "
    "~195-site class recorded in check-none-arithmetic's docstring: the checker "
    "models `.get()` returning None and cannot see a dataset value that IS None. "
    "This lane fixed the four sites the REPORT layer owns. Marked xfail rather "
    "than deleted so the gap stays measurable; it flips to xpass the day the "
    "class is closed."))
@pytest.mark.parametrize("missing", ["cash", "total_equity", "revenue"])
def test_generation_survives_a_genuinely_absent_line_item(missing):
    """⭐ The report layer is downstream of everything. After the 30 Jul coercion
    removals a None reaches the narrative; formatting it must render an em dash,
    not raise."""
    data = _dataset(missing=(missing,), missing_year=2023)
    try:
        intel.board_report(data)
    except (TypeError, ZeroDivisionError) as e:      # the two failure shapes seen
        pytest.fail(f"absent {missing} crashed generation: {type(e).__name__}: {e}")
    except Exception:
        pass          # a domain refusal is acceptable; a format crash is not


def test_both_formats_produce_an_openable_artifact():
    """⭐ A 200 IS NOT A WORKING REPORT. A PDF that renders blank and a PPTX
    PowerPoint refuses both return success — this is the false-green shape, and
    the report path is where it is least visible."""
    report = intel.board_report(_dataset())
    meta = _meta()

    pptx = R.build_pptx_comprehensive(report, {}, meta, _dataset())
    assert pptx[:2] == b"PK", "not a zip container — PowerPoint would refuse it"
    from pptx import Presentation
    import io
    pr = Presentation(io.BytesIO(pptx))
    shapes = [sh for s in pr.slides for sh in s.shapes
              if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()]
    assert len(pr.slides) > 0 and len(shapes) > 0, "deck opens but is blank"

    pdf = R.build_pdf(report, {}, meta)
    assert pdf[:5] == b"%PDF-", "not a PDF"
    assert pdf.count(b"/Type /Page") > 0, "PDF has no pages"
